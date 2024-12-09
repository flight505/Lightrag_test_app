import os
from zipfile import ZipFile
from xml.etree import ElementTree as etree
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

# helper functions

def docx_to_txt(input_dir, output_dir):
    """Converts .docx files in input_dir to .txt files in output_dir with tables formatted as Markdown
    and comments preserved."""
    
    for root, _, files in os.walk(input_dir):
        for file in files:
            if file.endswith(".docx"):
                docx_path = os.path.join(root, file)
                txt_path = os.path.join(output_dir, file[:-5] + ".txt")

                doc = Document(docx_path)
                full_text = []

                # Extract comments from .docx file
                comments = extract_comments_from_docx(docx_path)

                # Extract text from tables and format as Markdown
                for table in doc.tables:
                    table_rows = []
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            cell_text = []
                            for para in cell.paragraphs:
                                para_text_with_comments = process_paragraph(para, comments)
                                cell_text.append(para_text_with_comments)
                            row_text.append(" ".join(cell_text).strip())
                        table_rows.append(row_text)

                    # Add Markdown table formatting
                    if table_rows:
                        # Add header and separator if table exists
                        full_text.append("| " + " | ".join(table_rows[0]) + " |")
                        full_text.append("|" + " --- |" * len(table_rows[0]))

                        for row in table_rows[1:]:
                            full_text.append("| " + " | ".join(row) + " |")

                # Optionally, extract text from paragraphs outside of tables as well
                for para in doc.paragraphs:
                    para_text_with_comments = process_paragraph(para, comments)
                    full_text.append(para_text_with_comments)

                # Write the extracted text to a .txt file, keeping the structure
                with open(txt_path, "w", encoding="utf-8") as f:
                    f.write("\n".join(full_text))

def pptx_to_txt(input_dir, output_dir):
    """Converts .pptx files in input_dir to .txt files in output_dir with comments preserved."""
    
    for root, _, files in os.walk(input_dir):
        for file in files:
            if file.endswith(".pptx"):
                pptx_path = os.path.join(root, file)
                txt_path = os.path.join(output_dir, file[:-5] + ".txt")

                presentation = Presentation(pptx_path)
                full_text = []

                # Extract comments from .pptx file
                comments = extract_comments_from_pptx(pptx_path)
                
                # Extract text from slides and append comments
                for slide_idx, slide in enumerate(presentation.slides):
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            slide_text.append(shape.text)
                    
                    # Add comments for the slide
                    slide_comments = comments.get(slide_idx, [])
                    for comment in slide_comments:
                        slide_text.append(f"[Comment: {comment}]")
                    
                    full_text.append(f"Slide {slide_idx + 1}:\n" + "\n".join(slide_text))

                # Write the extracted text to a .txt file
                with open(txt_path, "w", encoding="utf-8") as f:
                    f.write("\n\n".join(full_text))

def xlsx_to_txt(input_dir, output_dir):
    """Converts .xlsx files in input_dir to .txt files in output_dir with tables formatted as Markdown
    and comments preserved."""
    
    for root, _, files in os.walk(input_dir):
        for file in files:
            if file.endswith(".xlsx"):
                xlsx_path = os.path.join(root, file)
                txt_path = os.path.join(output_dir, file[:-5] + ".txt")

                # Load excel workbook
                wb = load_workbook(xlsx_path, data_only=True)
                full_text = []

                # Iterate over worksheets
                for sheet in wb.worksheets:
                    sheet_text = [f"Sheet: {sheet.title}"]

                    # Collect the rows data
                    rows_data = []
                    for row in sheet.iter_rows(values_only=True):
                        row_data = [str(cell) if cell is not None else '' for cell in row]
                        rows_data.append(row_data)

                    # Ensure uniform row lengths
                    max_columns = max(len(row) for row in rows_data)
                    for row in rows_data:
                        if len(row) < max_columns:
                            row.extend([''] * (max_columns - len(row)))

                    # Add table header and Markdown separator if table exists
                    if rows_data:
                        sheet_text.append("| " + " | ".join(rows_data[0]) + " |")
                        sheet_text.append("|" + " --- |" * len(rows_data[0]))

                        for row in rows_data[1:]:
                            sheet_text.append("| " + " | ".join(row) + " |")

                    full_text.append("\n".join(sheet_text))

                # Write the extracted text to a .txt file, preserving the structure
                with open(txt_path, "w", encoding="utf-8") as f:
                    f.write("\n\n".join(full_text))

def extract_comments_from_docx(docx_path):
    """Extracts comments from the .docx file and returns them as a dictionary."""
    comments = {}
    
    with ZipFile(docx_path) as docx:
        try:
            comments_xml = docx.read('word/comments.xml')
            tree = etree.XML(comments_xml)
            comment_elements = tree.findall(".//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}comment")
            
            for comment in comment_elements:
                comment_id = comment.get("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}id")
                comment_text = ', '.join(comment.itertext()).strip()
                comments[comment_id] = comment_text
        except KeyError:
            print(f"No comments found in {docx_path}")

    return comments

def extract_comments_from_pptx(pptx_path):
    """Extracts comments from the .pptx file and returns them as a dictionary per slide."""
    comments = {}

    with ZipFile(pptx_path) as pptx:
        try:
            comments_xml = pptx.read('ppt/comments.xml')
            tree = etree.XML(comments_xml)
            comment_elements = tree.findall(".//{http://schemas.openxmlformats.org/presentationml/2006/main}cm")
            
            for comment in comment_elements:
                slide_idx = int(comment.get("{http://schemas.openxmlformats.org/presentationml/2006/main}slideId")) - 1
                comment_text = ', '.join(comment.itertext()).strip()
                
                if slide_idx not in comments:
                    comments[slide_idx] = []
                comments[slide_idx].append(comment_text)
        except KeyError:
            print(f"No comments found in {pptx_path}")

    return comments

def extract_comments_from_xlsx(xlsx_path):
    """Extracts comments from the .xlsx file and returns them as a dictionary with sheet and cell as keys."""
    comments = {}

    with ZipFile(xlsx_path) as xlsx:
        try:
            comments_xml = xlsx.read('xl/comments1.xml')
            tree = etree.XML(comments_xml)
            comment_elements = tree.findall(".//{http://schemas.openxmlformats.org/spreadsheetml/2006/main}comment")
            
            for comment in comment_elements:
                ref = comment.get("ref")
                sheet_name = comment.get("sheet")
                comment_text = ', '.join(comment.itertext()).strip()
                
                comments[(sheet_name, ref)] = comment_text
        except KeyError:
            print(f"No comments found in {xlsx_path}")

    return comments

def process_paragraph(paragraph, comments):
    """Processes a paragraph, inserting comments after the relevant text snippets."""
    para_element = paragraph._element
    namespace_uri = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
    
    etree.register_namespace('w', namespace_uri)
    comment_start_elements = para_element.findall(".//w:commentRangeStart", namespaces={'w': namespace_uri})
    comment_end_elements = para_element.findall(".//w:commentRangeEnd", namespaces={'w': namespace_uri})

    comment_ranges = {}
    for start_elem in comment_start_elements:
        comment_id = start_elem.get(qn("w:id"))
        if comment_id not in comment_ranges:
            comment_ranges[comment_id] = {"start": start_elem, "end": None}
        
    for end_elem in comment_end_elements:
        comment_id = end_elem.get(qn("w:id"))
        if comment_id in comment_ranges:
            comment_ranges[comment_id]["end"] = end_elem

    snippets = []
    text_elements = para_element.findall(".//w:t", namespaces={'w': namespace_uri})
    for elem in text_elements:
        snippets.append(elem.text or '')

    snippet_comments = ["" for _ in range(len(snippets))]
    
    for comment_id, range_info in comment_ranges.items():
        start_elem = range_info["start"]
        end_elem = range_info["end"]
        comment_text = comments.get(comment_id, "")
        
        start_index = para_element.index(start_elem)
        end_index = para_element.index(end_elem)
        
        if 0 <= start_index < len(snippets):
            snippet_comments[start_index] += f" [Comment: {comment_text}]"
        if 0 <= end_index < len(snippets) and end_index != start_index:
            snippet_comments[end_index] += f" [Comment: {comment_text}]"

    para_text_with_comments = ''.join(f"{snippets[i]}{snippet_comments[i]}" for i in range(len(snippets)))
    
    return para_text_with_comments

# main function

def convert_files(input_dir, output_dir):
    """Converts .docx, .pptx, and .xlsx files in input_dir to .txt files in output_dir with Markdown formatting
    and comments preserved."""
    docx_to_txt(input_dir, output_dir)
    pptx_to_txt(input_dir, output_dir)
    xlsx_to_txt(input_dir, output_dir)

# declare paths
input_dirs=[
    #input path
]
output_dir= # output path

# run
for input_dir in input_dirs:
    if __name__ == "__main__":
        convert_files(input_dir, output_dir)